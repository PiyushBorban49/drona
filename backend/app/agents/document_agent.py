"""
Dronacharya v3 — Document Ingestion Agent
Extracts text from PDF, DOCX, and PPTX files for AI synthesis.
"""
import os
import pypdf
import docx
from pptx import Presentation
from typing import Dict
from langchain_groq import ChatGroq
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import JsonOutputParser
from app.config import get_settings

async def extract_text_pypdf(file_path: str) -> str:
    """Extract text from PDF using pypdf."""
    text = ""
    try:
        reader = pypdf.PdfReader(file_path)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    except Exception as e:
        print(f"PDF extraction error: {e}")
    return text

async def extract_text_docx(file_path: str) -> str:
    """Extract text from DOCX using python-docx."""
    text = ""
    try:
        doc = docx.Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        print(f"DOCX extraction error: {e}")
    return text

async def extract_text_pptx(file_path: str) -> str:
    """Extract text from PPTX using python-pptx."""
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"PPTX extraction error: {e}")
    return text

async def synthesize_document_content(filename: str, text: str) -> Dict:
    """Use AI to synthesize document text into structured knowledge."""
    if not text.strip():
        return {"summary": "Document extraction failed.", "concepts": [], "key_points": []}

    settings = get_settings()
    llm = ChatGroq(
        model_name="llama-3.3-70b-versatile", 
        temperature=0.1,
        api_key=settings.GROQ_API_KEY
    )
    
    prompt = ChatPromptTemplate.from_messages([
        ("system", "You are an expert researcher. Synthesize the following raw text from a document into a concise overview, key concepts, and main takeaways. Return valid JSON."),
        ("user", "Document: {filename}\nContent:\n{text}\n\nReturn JSON: {{ 'summary': '...', 'concepts': [{{ 'title': '...', 'description': '...' }}], 'key_points': ['...', '...'] }}")
    ])

    chain = prompt | llm | JsonOutputParser()
    
    try:
        # Cap text at 15000 chars for LLM context
        return await chain.ainvoke({"filename": filename, "text": text[:15000]})
    except Exception as e:
        print(f"Document synthesis error: {e}")
        return {"summary": "Document analysis failed.", "concepts": []}

async def process_document(file_path: str) -> Dict:
    """Main orchestration for Document ingestion."""
    ext = os.path.splitext(file_path)[1].lower()
    text = ""
    
    if ext == '.pdf':
        text = await extract_text_pypdf(file_path)
    elif ext == '.docx':
        text = await extract_text_docx(file_path)
    elif ext == '.pptx':
        text = await extract_text_pptx(file_path)
    else:
        return {"success": False, "error": f"Unsupported file type: {ext}"}

    if not text.strip():
        return {"success": False, "error": "No text could be extracted from the document."}

    analysis = await synthesize_document_content(os.path.basename(file_path), text)
    
    return {
        "success": True,
        "filename": os.path.basename(file_path),
        "analysis": analysis,
        "raw_text_preview": text[:500]
    }
